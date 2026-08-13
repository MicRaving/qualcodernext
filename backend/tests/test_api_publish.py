"""API tests for the Smart Publisher (docx/pptx/xlsx export) endpoints.

Synthetic project: two files, three codes (one with a memo), text codings
with memos. The publish router is not wired into ``api/v1/router.py`` yet —
the app supervisor includes it — so these tests mount both the v1 router and
the publish router on a fresh app and run regardless of the wiring state.
"""

from __future__ import annotations

from io import BytesIO

import pytest
from docx import Document
from fastapi import FastAPI
from httpx import ASGITransport, AsyncClient
from openpyxl import load_workbook
from pptx import Presentation

from qualcoder_api.api.v1.publish import router as publish_router
from qualcoder_api.api.v1.router import router as v1_router
from qualcoder_api.services import publish_service

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

publish_app = FastAPI()
publish_app.include_router(v1_router, prefix="/api/v1")
publish_app.include_router(publish_router, prefix="/api/v1")


@pytest.fixture(autouse=True)
def isolated_settings(tmp_path, monkeypatch):
    """Keep the developer's real user settings out of the run."""
    from qualcoder_api.services import user_settings

    monkeypatch.setattr(user_settings, "SETTINGS_FILE", tmp_path / "settings.json")


@pytest.fixture
async def publish_client(tmp_path):
    transport = ASGITransport(app=publish_app)
    async with AsyncClient(transport=transport, base_url="http://test") as c:
        target = tmp_path / "publish.qda"
        res = await c.post(
            "/api/v1/projects",
            json={"project_path": str(target), "codername": "tester"},
        )
        assert res.status_code == 200, res.text
        yield c
        await c.post("/api/v1/projects/close")


@pytest.fixture
async def publish_dataset(publish_client):
    """Files a/b, codes C1..C3 (C1 with a memo), text codings with memos."""
    client = publish_client

    async def import_file(name, text):
        res = await client.post(
            "/api/v1/sources/import",
            files={"file": (name, text, "text/plain")},
        )
        assert res.status_code == 200, res.text
        return res.json()["id"]

    f1 = await import_file("a.txt", "alpha beta\ngamma delta")
    f2 = await import_file("b.txt", "epsilon zeta\neta theta")

    cids = {}
    for name in ("C1", "C2", "C3"):
        res = await client.post("/api/v1/codes", json={"name": name})
        assert res.status_code == 201, res.text
        cids[name] = res.json()["cid"]
    res = await client.patch(f"/api/v1/codes/{cids['C1']}", json={"memo": "about themes"})
    assert res.status_code == 200, res.text

    async def code(fid, name, text, pos0, pos1, memo=""):
        res = await client.post(
            "/api/v1/codings/text",
            json={
                "cid": cids[name], "fid": fid, "seltext": text,
                "pos0": pos0, "pos1": pos1, "memo": memo,
            },
        )
        assert res.status_code == 201, res.text

    # f1: C1 twice (two memos) + C2; f2: C1 + C3.
    await code(f1, "C1", "alpha", 0, 5, "first")
    await code(f1, "C1", "beta", 6, 10, "extra")
    await code(f1, "C2", "gamma", 11, 16, "second")
    await code(f2, "C1", "epsilon", 0, 7, "third")
    await code(f2, "C3", "zeta", 8, 12, "fourth")

    return {"f1": f1, "f2": f2, "cids": cids}


# ----------------------------------------------------------------------
# Builder round-trips
# ----------------------------------------------------------------------


def test_build_docx_roundtrip():
    sections = [
        {"paragraph": "intro text"},
        {"heading": "Coding counts"},
        {"table": [["Code", "Count"], ["ThemeA", 2], ["ThemeB", 1]]},
        {"quote": "[a.txt] alpha"},
    ]
    data = publish_service.build_docx("Test report", sections)
    doc = Document(BytesIO(data))
    assert doc.paragraphs[0].text == "Test report"
    texts = [p.text for p in doc.paragraphs]
    assert "intro text" in texts
    assert "Coding counts" in texts
    assert "[a.txt] alpha" in texts
    assert len(doc.tables) == 1
    table = doc.tables[0]
    assert [cell.text for cell in table.rows[0].cells] == ["Code", "Count"]
    assert table.cell(1, 0).text == "ThemeA"
    assert table.cell(1, 1).text == "2"
    assert table.rows[0].cells[0].paragraphs[0].runs[0].font.bold is True


def test_build_docx_empty_table():
    data = publish_service.build_docx("Empty", [{"table": []}])
    doc = Document(BytesIO(data))
    assert doc.paragraphs[0].text == "Empty"
    assert len(doc.tables) == 0


def test_build_pptx_roundtrip():
    slides = [
        {"title": "ThemeA", "bullets": ["a.txt — alpha", "b.txt — beta"], "memo": "memo A"},
        {"title": "ThemeB", "bullets": [], "memo": ""},
    ]
    data = publish_service.build_pptx("Deck", slides)
    prs = Presentation(BytesIO(data))
    assert len(prs.slides) == 3  # title slide + one per code
    assert prs.slides[0].shapes.title.text == "Deck"
    assert prs.slides[1].shapes.title.text == "ThemeA"
    bullets = [p.text for p in prs.slides[1].placeholders[1].text_frame.paragraphs]
    assert bullets[0] == "a.txt — alpha"
    assert bullets[1] == "b.txt — beta"
    assert bullets[-1] == "Memo: memo A"
    # Empty-bullet slide renders a placeholder line, not a blank slide.
    empty = [p.text for p in prs.slides[2].placeholders[1].text_frame.paragraphs]
    assert empty == ["—"]


def test_build_xlsx_roundtrip():
    sheets = [
        {"name": "Code frequencies", "headers": ["Code", "Count"], "rows": [["ThemeA", 2]]},
        {"name": "Empty sheet", "headers": [], "rows": []},
    ]
    data = publish_service.build_xlsx(sheets)
    wb = load_workbook(BytesIO(data))
    assert wb.sheetnames == ["Code frequencies", "Empty sheet"]
    ws = wb["Code frequencies"]
    assert ws["A1"].value == "Code"
    assert ws["A1"].font.bold is True
    assert ws["A2"].value == "ThemeA"
    assert ws["B2"].value == 2  # numeric values survive as numbers
    assert wb["Empty sheet"].max_row == 1  # header-only sheet still exists


# ----------------------------------------------------------------------
# Endpoint behaviour
# ----------------------------------------------------------------------


async def test_publish_code_frequencies_docx(publish_client, publish_dataset):
    client = publish_client
    res = await client.post(
        "/api/v1/publish/from-report",
        json={"report": "code-frequencies", "format": "docx"},
    )
    assert res.status_code == 200, res.text
    assert res.headers["content-type"].startswith(DOCX_MIME)
    assert "attachment" in res.headers["content-disposition"]
    assert "code-frequencies.docx" in res.headers["content-disposition"]
    assert len(res.content) > 1000
    doc = Document(BytesIO(res.content))
    assert doc.paragraphs[0].text == "Code frequencies"
    assert len(doc.tables) == 1
    table = doc.tables[0]
    assert [cell.text for cell in table.rows[0].cells] == ["Code", "Category", "Segments"]
    body = [[cell.text for cell in row.cells] for row in table.rows[1:]]
    # C1 (3 segments) ranks first: [name, category, count].
    assert body[0][0] == "C1"
    assert body[0][2] == "3"


async def test_publish_code_segments_xlsx(publish_client, publish_dataset):
    client = publish_client
    res = await client.post(
        "/api/v1/publish/from-report",
        json={"report": "code-segments", "format": "xlsx"},
    )
    assert res.status_code == 200, res.text
    assert res.headers["content-type"].startswith(XLSX_MIME)
    wb = load_workbook(BytesIO(res.content))
    assert wb.sheetnames == ["Code segments"]
    ws = wb["Code segments"]
    assert [c.value for c in ws[1]] == ["Code", "Category", "File", "Segment", "Owner", "Date"]
    rows = list(ws.iter_rows(min_row=2, values_only=True))
    assert len(rows) == 5
    assert all(r[2] in ("a.txt", "b.txt") for r in rows)
    assert any(
        r[0] == "C1" and r[2] == "a.txt" and r[3] == "alpha" and r[4] == "default"
        for r in rows
    )


async def test_publish_summary_table_xlsx_and_docx(publish_client, publish_dataset):
    client = publish_client
    res = await client.post(
        "/api/v1/publish/from-report",
        json={"report": "summary-table", "format": "xlsx"},
    )
    assert res.status_code == 200, res.text
    wb = load_workbook(BytesIO(res.content))
    ws = wb["Summary table"]
    headers = [c.value for c in ws[1]]
    assert headers[0] == "Document"
    assert headers[1:] == ["C1", "C2", "C3"]
    rows = {row[0]: row for row in ws.iter_rows(min_row=2, values_only=True)}
    assert rows["a.txt"][1] == "first — extra"
    assert rows["b.txt"][2] in (None, "")

    res = await client.post(
        "/api/v1/publish/from-report",
        json={"report": "summary-table", "format": "docx", "options": {"scope": "case"}},
    )
    assert res.status_code == 200, res.text
    doc = Document(BytesIO(res.content))
    assert doc.paragraphs[0].text == "Summary table"
    table = doc.tables[0]
    assert table.rows[0].cells[0].text == "Case"


async def test_publish_coder_comparison_docx(publish_client, publish_dataset):
    client = publish_client
    res = await client.post(
        "/api/v1/publish/from-report",
        json={"report": "coder-comparison", "format": "docx"},
    )
    assert res.status_code == 200, res.text
    doc = Document(BytesIO(res.content))
    assert doc.paragraphs[0].text == "Coder comparison"
    table = doc.tables[0]
    assert [cell.text for cell in table.rows[0].cells] == ["Coder", "Codings", "Files"]
    assert table.cell(1, 0).text == "default"
    assert table.cell(1, 1).text == "5"


async def test_publish_codebook_docx(publish_client, publish_dataset):
    client = publish_client
    res = await client.post(
        "/api/v1/publish/from-report",
        json={"report": "codebook", "format": "docx"},
    )
    assert res.status_code == 200, res.text
    doc = Document(BytesIO(res.content))
    texts = "\n".join(p.text for p in doc.paragraphs)
    assert texts.startswith("Codebook")
    assert "C1" in texts
    assert "about themes" in texts  # code memo is included


async def test_publish_code_segments_pptx(publish_client, publish_dataset):
    client = publish_client
    res = await client.post(
        "/api/v1/publish/from-report",
        json={"report": "code-segments", "format": "pptx"},
    )
    assert res.status_code == 200, res.text
    assert res.headers["content-type"].startswith(PPTX_MIME)
    prs = Presentation(BytesIO(res.content))
    titles = [slide.shapes.title.text for slide in prs.slides]
    assert titles[0] == "Codes by segments"
    assert set(titles[1:]) == {"C1", "C2", "C3"}
    c1 = next(i for i, t in enumerate(titles) if t == "C1")
    bullets = [p.text for p in prs.slides[c1].placeholders[1].text_frame.paragraphs]
    assert any("alpha" in b for b in bullets)
    assert bullets[-1] == "Memo: about themes"


async def test_publish_code_frequencies_pptx_top_n(publish_client, publish_dataset):
    client = publish_client
    res = await client.post(
        "/api/v1/publish/from-report",
        json={"report": "code-frequencies", "format": "pptx", "options": {"top_n": 1}},
    )
    assert res.status_code == 200, res.text
    prs = Presentation(BytesIO(res.content))
    assert len(prs.slides) == 2  # title + top code only
    assert prs.slides[1].shapes.title.text == "C1"
    bullets = [p.text for p in prs.slides[1].placeholders[1].text_frame.paragraphs]
    assert bullets[1] == "Segments: 3"


async def test_publish_unknown_report_and_format(publish_client, publish_dataset):
    client = publish_client
    res = await client.post(
        "/api/v1/publish/from-report",
        json={"report": "bogus", "format": "docx"},
    )
    assert res.status_code == 422
    assert "unknown report: bogus" in res.json()["detail"]

    res = await client.post(
        "/api/v1/publish/from-report",
        json={"report": "code-frequencies", "format": "pdf"},
    )
    assert res.status_code == 422
    assert "unknown format: pdf" in res.json()["detail"]


async def test_publish_pptx_unsupported_reports(publish_client, publish_dataset):
    client = publish_client
    for report in ("coder-comparison", "codebook", "summary-table"):
        res = await client.post(
            "/api/v1/publish/from-report",
            json={"report": report, "format": "pptx"},
        )
        assert res.status_code == 422
        assert "PowerPoint is only supported" in res.json()["detail"]


async def test_publish_summary_table_bad_scope(publish_client, publish_dataset):
    client = publish_client
    res = await client.post(
        "/api/v1/publish/from-report",
        json={"report": "summary-table", "format": "docx", "options": {"scope": "bogus"}},
    )
    assert res.status_code == 422
    assert "options.scope" in res.json()["detail"]


async def test_publish_requires_open_project(tmp_path):
    transport = ASGITransport(app=publish_app)
    async with AsyncClient(transport=transport, base_url="http://test") as c:
        res = await c.post(
            "/api/v1/publish/from-report",
            json={"report": "code-frequencies", "format": "docx"},
        )
        assert res.status_code == 409
