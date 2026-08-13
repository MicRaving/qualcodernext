"""Smart Publisher — render report datasets into Office documents.

Pure renderers: ``build_docx`` / ``build_pptx`` / ``build_xlsx`` take plain
data (title, sections, slides, sheets) and return the serialized file bytes.
The API layer (``api/v1/publish.py``) pulls the report data through
``report_service`` and maps it onto these inputs. No images in v1.
"""

from __future__ import annotations

from datetime import date
from io import BytesIO
from typing import Any

from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Pt
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation

_HEADER_FILL = "D9E2F3"
_XLSX_HEADER_FILL = "4472C4"


def _text(value: Any) -> str:
    """Render any cell value as display text."""
    return "" if value is None else str(value)


# ----------------------------------------------------------------------
# Word
# ----------------------------------------------------------------------


def _shade_docx_cells(cells: list[Any], fill: str) -> None:
    for cell in cells:
        tc_pr = cell._tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:fill"), fill)
        tc_pr.append(shd)


def _add_docx_table(doc: Any, rows: list[list[str]]) -> None:
    if not rows:
        return
    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
    table.style = "Table Grid"
    for ri, row in enumerate(rows):
        for ci, value in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = value
            for run in cell.paragraphs[0].runs:
                run.font.size = Pt(10)
                if ri == 0:
                    run.font.bold = True
    _shade_docx_cells(list(table.rows[0].cells), _HEADER_FILL)


def build_docx(title: str, sections: list[dict[str, Any]]) -> bytes:
    """Render a Word document from a block list.

    Each section supports: ``heading`` (level-1 heading), ``paragraph``,
    ``quote`` (indented italic) and ``table`` (list-of-lists; the first row
    is the header). A section may carry several of these keys.
    """
    doc = Document()
    doc.add_heading(_text(title), level=0)
    for section in sections:
        if "heading" in section:
            doc.add_heading(_text(section["heading"]), level=1)
        if "paragraph" in section:
            doc.add_paragraph(_text(section["paragraph"]))
        if "quote" in section:
            paragraph = doc.add_paragraph()
            paragraph.paragraph_format.left_indent = Pt(18)
            run = paragraph.add_run(_text(section["quote"]))
            run.italic = True
        table = section.get("table")
        if table:
            _add_docx_table(doc, [[_text(cell) for cell in row] for row in table])
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ----------------------------------------------------------------------
# PowerPoint
# ----------------------------------------------------------------------


def build_pptx(title: str, slides: list[dict[str, Any]]) -> bytes:
    """Render a PowerPoint deck: title slide + one content slide per item.

    Each slide entry carries ``title``, ``bullets`` (list of lines) and an
    optional ``memo`` rendered as a trailing ``Memo: …`` line.
    """
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = _text(title)
    title_slide.placeholders[1].text = f"Generated {date.today().isoformat()}"

    for item in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = _text(item.get("title", ""))
        lines = [_text(line) for line in item.get("bullets", [])]
        memo = item.get("memo")
        if memo:
            lines.append(f"Memo: {_text(memo)}")
        if not lines:
            lines = ["—"]
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, line in enumerate(lines):
            paragraph = body.paragraphs[0] if i == 0 else body.add_paragraph()
            paragraph.text = line
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ----------------------------------------------------------------------
# Excel
# ----------------------------------------------------------------------


def build_xlsx(sheets: list[dict[str, Any]]) -> bytes:
    """Render an Excel workbook: one styled sheet per entry.

    Each sheet carries ``name``, ``headers`` and ``rows``; the header row is
    bold white on a colored fill, column widths are sized to the content.
    """
    wb = Workbook()
    wb.remove(wb.active)
    for sheet in sheets:
        name = _text(sheet.get("name", "Sheet"))[:31]
        headers = [_text(header) for header in sheet.get("headers", [])]
        rows = list(sheet.get("rows", []))
        ws = wb.create_sheet(title=name)

        ncols = len(headers)
        if rows:
            ncols = max(ncols, max(len(row) for row in rows))
        if not headers:
            headers = [f"Column {i + 1}" for i in range(ncols)]

        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill("solid", fgColor=_XLSX_HEADER_FILL)
        wrap = Alignment(vertical="top", wrap_text=True)
        for ci, header in enumerate(headers[:ncols], start=1):
            cell = ws.cell(row=1, column=ci, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = wrap
        for ri, row in enumerate(rows, start=2):
            for ci in range(ncols):
                value = row[ci] if ci < len(row) else ""
                ws.cell(row=ri, column=ci + 1, value=value)

        for ci in range(ncols):
            longest = len(headers[ci])
            for row in rows:
                if ci < len(row):
                    longest = max(longest, len(_text(row[ci])))
            ws.column_dimensions[get_column_letter(ci + 1)].width = min(
                40, max(12, longest + 2)
            )
        ws.freeze_panes = "A2"
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()
